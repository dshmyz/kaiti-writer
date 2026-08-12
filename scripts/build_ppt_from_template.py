#!/usr/bin/env python3
"""基于内置北航答辩 PPT 模板生成开题汇报 PPT。

策略：不新建幻灯片，只从模板里**挑选**需要的页、替换其占位文字、删除未用页，
保留模板全部母版/版式/配色/图形（同 build_from_template.py 的"改而不建"思路）。
"""
from __future__ import annotations
import argparse, json, re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    raise SystemExit(
        "❌ 缺少 python-pptx，无法生成 PPT。\n"
        "  安装：pip install python-pptx（或 python3 -m pip install python-pptx）\n"
        "  不想安装：可改为只要「markdown 逐页大纲」，自己粘进 PPT。"
    )

# 各模板的页角色索引（1-based，来自实际探查）
TEMPLATE_MAP = {
    "模板1": {"cover": 1, "toc": 3, "section": [5, 10, 15, 20], "content": [7, 11, 12], "thanks": 38},
    "模板2": {"cover": 1, "toc": 3, "section": [4, 8, 12, 16, 20, 24], "content": [9, 10, 14], "thanks": 32},
    "模板3": {"cover": 1, "toc": 2, "section": [3, 9, 16, 22, 27], "content": [5, 7, 12], "thanks": 36},
}

# 模板作者水印 / 样例信息，一律清掉或替换
WATERMARKS = ["星空情报站", "微信公众号", "公众号", "微信", "BEIHANG UNIVERSITY",
              "北京航空航天大学", "Beihang University"]

# 判定为"占位文字"的特征（含则视为可替换的占位）
PLACEHOLDER_HINTS = [
    "请输入", "请在此输入", "点击此处", "点击输入", "点击这里", "在这里输入", "在此录入",
    "单击此处", "ADD YOUR TITLE", "SAMPLE TITLE", "SIMPLE TITLE", "JUNE 12th",
    "你的标题", "你的模板你做主", "右键点击图片", "输入你的", "输入小标题", "输入标题",
    "输入文字说明", "输入你的观点", "你的模板", "Copy paste fonts", "Key words",
    "做的好的地方", "输入分点", "输入主题", "在这里输入", "输入你的内容",
    "BEIHANG UNIVERSITY", "北京航空航天大学", "摘 要",
]

# 纯装饰性文字：清空但不占内容位
DECOR_HINTS = ["JUNE 12th", "做的好的地方", "需要继续提升的地方", "/" * 6,
               "1997.11.12", "XXXX-XX", "Common PPT Template", "GRADUATION REPLY",
               "General template for teaching", "ADD YOUR TITLE", "请输入你的副标题",
               "请输入你的标题说明", "以简洁为主说明本章", "摘 要", "关键词0",
               "Welcoming", "Thank you for your", "图示页排版", "纯图片页排版",
               "小标题一", "小标题二", "小标题三", "小标题四", "段落标题",
               "BEIHANG UNIVERSITY", "北京航空航天大学"]

# 纯装饰的"占位符号"（整块就是这些字符时才清），如 XX、XX%、A/B/C 角标、示例日期
DECOR_EXACT = re.compile(
    r"^(X{1,4}%?|[A-F]|VS|\d{1,2}%|20\d{2}|\d{4}\.\d{2}\.\d{2}"
    r"|20\d{2}\s*年\s*\d{1,2}\s*月\s*\d{1,2}\s*日)$")





def is_placeholder(t: str) -> bool:
    t = t.strip()
    return bool(t) and any(h in t for h in PLACEHOLDER_HINTS)


def set_text_keep_style(shape, text: str):
    """只改第一个 run 的文字、删掉其余 run，保留字体/字号/颜色等全部样式。"""
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p0.text = text
    # 删掉后续段落，避免残留占位行
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)


def estimate_lines(text: str, box_width_emu: int, font_pt: float) -> int:
    """估算文本在给定宽度的文本框里需要多少行。"""
    # 中文字符约 font_pt * 0.55mm 宽，英文约 font_pt * 0.3mm；取混合均值 0.45mm
    char_width_mm = font_pt * 0.45
    box_width_mm = box_width_emu / 36000  # EMU → mm
    chars_per_line = max(1, int(box_width_mm / char_width_mm))
    # 按换行符和自动换行估算行数
    total_lines = 0
    for paragraph in text.split("\n"):
        if not paragraph.strip():
            total_lines += 1
        else:
            total_lines += max(1, -(-len(paragraph) // chars_per_line))  # 向上取整
    return total_lines


def fit_text_to_shape(shape, text: str, min_font_pt: float = 9.0):
    """将文本写入形状，若溢出则逐步缩小字号直到能装下；保留原始样式为起点。"""
    from pptx.util import Pt
    tf = shape.text_frame
    # 估算当前字号
    font_pt = 12.0
    if tf.paragraphs and tf.paragraphs[0].runs:
        rpr = tf.paragraphs[0].runs[0].font
        if rpr.size:
            font_pt = rpr.size.pt
    # 估算可用行数
    available_lines = max(1, int(shape.height / (font_pt * 12700 * 1.3)))  # 1.3x line spacing
    needed = estimate_lines(text, shape.width, font_pt)
    # 如果溢出，逐步缩字号
    while needed > available_lines and font_pt > min_font_pt:
        font_pt -= 0.5
        needed = estimate_lines(text, shape.width, font_pt)
    # 写入文字
    set_text_keep_style(shape, text)
    # 如果缩了字号，应用到第一个 run
    if tf.paragraphs and tf.paragraphs[0].runs:
        final_size = Pt(font_pt)
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.size = final_size
    if needed > available_lines:
        print(f"  ⚠ 文本框装不下，已缩到 {font_pt:.1f}pt 仍可能溢出: \"{text[:30]}…\"")


def iter_text_shapes(container):
    """递归遍历（含 GROUP 组合形状内部）所有有文字的形状。

    模板1/模板2 的封面与多数页把文字放在组合形状里，只看顶层 shape 会全部漏掉。
    """
    for sh in container.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_text_shapes(sh)
        elif sh.has_text_frame and sh.text_frame.text.strip():
            yield sh


def text_shapes(slide):
    """按视觉顺序（上→下、左→右）返回有文字的形状（含组合内部）。"""
    return sorted(iter_text_shapes(slide), key=lambda s: (s.top or 0, s.left or 0))


def clean_watermarks(slide):
    """清掉模板作者水印（把含水印的整块文字替换成空）。含组合形状内部。"""
    for sh in iter_text_shapes(slide):
        t = sh.text_frame.text
        if any(w in t for w in WATERMARKS):
            cleaned = t
            for w in WATERMARKS:
                cleaned = cleaned.replace(w, "")
            cleaned = re.sub(r"[：:）)\s]+$", "", cleaned).strip()
            set_text_keep_style(sh, cleaned)


def clear_decor(slide):
    """清空纯装饰性的模板样例文字（英文通用模板说明、装饰串、示例日期、XX 类占位符号）。"""
    for sh in iter_text_shapes(slide):
        t = sh.text_frame.text
        if any(d in t for d in DECOR_HINTS) or DECOR_EXACT.match(t.strip()):
            set_text_keep_style(sh, "")


def fill_cover(slide, data):
    """封面：最长的那块文字当题目，含"汇报人/答辩人/指导老师/日期"的块按标签填。"""
    cover = data.get("cover", {})
    title = data.get("title", "")
    shapes = text_shapes(slide)
    # 题目 = 字数最多且非标签的那块
    cand = [sh for sh in shapes
            if not any(k in sh.text_frame.text for k in ["汇报人", "答辩人", "指导老师", "年", "BEIHANG"])]
    if cand and title:
        target = max(cand, key=lambda s: len(s.text_frame.text))
        set_text_keep_style(target, title)
    for sh in shapes:
        t = sh.text_frame.text
        if ("汇报人" in t or "答辩人" in t) and cover.get("作者姓名"):
            lbl = "答辩人" if "答辩人" in t else "汇报人"
            set_text_keep_style(sh, f"{lbl}：{cover['作者姓名']}")
        elif "指导老师" in t and cover.get("指导教师"):
            set_text_keep_style(sh, f"指导老师：{cover['指导教师']}")
        elif re.search(r"20\d{2}\s*年", t) and cover.get("日期"):
            set_text_keep_style(sh, cover["日期"])
    clear_decor(slide)
    clean_watermarks(slide)


def _group_overlapping_shapes(shapes, tol=50000):
    """将位置相近的形状归为一组（容忍 50000 EMU ≈ 0.5mm）。

    模板的目录页常把多组视觉变体叠在相同位置，需要先分组再每组只取一个。
    """
    groups = []
    used = set()
    shapes_list = list(shapes)
    for i, a in enumerate(shapes_list):
        if i in used:
            continue
        group = [a]
        used.add(i)
        for j, b in enumerate(shapes_list):
            if j in used:
                continue
            if (abs((a.top or 0) - (b.top or 0)) < tol and
                    abs((a.left or 0) - (b.left or 0)) < tol):
                group.append(b)
                used.add(j)
        groups.append(group)
    return groups


def fill_toc(slide, chapters):
    """目录页：检测重叠形状组，每组只用面积最大的一个，按视觉位置排序后填入序号和标题。"""
    all_shapes = list(text_shapes(slide))
    # 分离数字形状和标题形状
    num_shapes = [s for s in all_shapes if re.fullmatch(r"0?\d{1,2}", s.text_frame.text.strip())]
    title_shapes = [s for s in all_shapes if is_placeholder(s.text_frame.text) and "ADD YOUR" not in s.text_frame.text]

    # 对数字形状分组（检测重叠），每组取面积最大的
    num_groups = _group_overlapping_shapes(num_shapes)
    num_slots = [max(g, key=lambda s: (s.width or 0) * (s.height or 0)) for g in num_groups]
    num_slots.sort(key=lambda s: (s.top or 0, s.left or 0))

    # 对标题形状分组（检测重叠），每组取面积最大的
    title_groups = _group_overlapping_shapes(title_shapes)
    title_slots = [max(g, key=lambda s: (s.width or 0) * (s.height or 0)) for g in title_groups]
    title_slots.sort(key=lambda s: (s.top or 0, s.left or 0))

    # 填入标题
    for i, sh in enumerate(title_slots):
        set_text_keep_style(sh, chapters[i] if i < len(chapters) else "")
    # 填入序号
    for i, sh in enumerate(num_slots):
        set_text_keep_style(sh, f"{i + 1:02d}" if i < len(chapters) else "")
    # 多余的清空
    for i in range(len(chapters), len(title_slots)):
        set_text_keep_style(title_slots[i], "")
    for i in range(len(chapters), len(num_slots)):
        set_text_keep_style(num_slots[i], "")
    clear_decor(slide)
    clean_watermarks(slide)


def fill_section(slide, idx, name):
    """章节过渡页：序号写 01/02…，第一个标题占位写章节名，其余占位清空。"""
    first = True
    for sh in text_shapes(slide):
        t = sh.text_frame.text.strip()
        if re.fullmatch(r"0?\d{1,2}", t):
            set_text_keep_style(sh, f"{idx:02d}")
        elif is_placeholder(t):
            if first and "ADD YOUR" not in t:
                set_text_keep_style(sh, name)
                first = False
            else:
                set_text_keep_style(sh, "")
    clear_decor(slide)
    clean_watermarks(slide)


def fill_content(slide, title, bullets, layout="text_only", extra=None):
    """内容页：根据 layout 类型填充不同内容。

    layout: text_only / image_center / chart / table / big_number / comparison / timeline
    extra: 额外数据（图片路径、图表数据、表格数据等）
    """
    clear_decor(slide)

    if layout == "image_center" and extra:
        _fill_image_slide(slide, title, extra)
        return
    if layout == "chart" and extra:
        _fill_chart_slide(slide, title, extra)
        return
    if layout == "table" and extra:
        _fill_table_slide(slide, title, extra)
        return
    if layout == "big_number" and extra:
        _fill_big_number_slide(slide, title, extra)
        return
    if layout == "comparison" and extra:
        _fill_comparison_slide(slide, title, extra)
        return
    if layout == "timeline" and extra:
        _fill_timeline_slide(slide, title, extra)
        return

    # 默认：文字内容页
    slots = [sh for sh in text_shapes(slide) if is_placeholder(sh.text_frame.text)]
    if not slots:
        clean_watermarks(slide)
        return

    # 改进：结合位置和内容提示选形状，不再只按面积
    slide_height = slide.slide_height or 6858000  # 默认高度
    title_shape, content_shape = _select_title_and_content(slots, slide_height)

    # 写标题
    set_text_keep_style(title_shape, title)
    # 写内容（多条 bullets，用 fit 防溢出）
    full_text = "\n".join(bullets)
    fit_text_to_shape(content_shape, full_text)
    # 剩余小占位块清空
    used = {content_shape, title_shape}
    for sh in slots:
        if sh not in used:
            set_text_keep_style(sh, "")
    clean_watermarks(slide)


def set_slide_notes(slide, notes_text: str):
    """设置幻灯片的演讲者备注。"""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def _select_title_and_content(slots, slide_height):
    """从候选形状中选出标题和内容框。

    策略：
    1. 顶部区域（top < 30%）的形状优先作为标题
    2. 包含标题提示的形状优先作为标题
    3. 面积最大的非标题形状作为内容
    """
    # 标记每个形状的"标题得分"
    scored = []
    for s in slots:
        score = 0
        top_ratio = (s.top or 0) / slide_height
        text = s.text_frame.text.lower()
        # 顶部区域加分
        if top_ratio < 0.3:
            score += 10
        elif top_ratio < 0.5:
            score += 5
        # 标题提示加分
        if any(h in text for h in ["标题", "title", "请输入", "单击此处"]):
            score += 20
        # 面积作为基础分（越大越可能是内容框）
        area = (s.width or 0) * (s.height or 0)
        score += min(area // 1000000000, 5)  # 限制面积权重
        scored.append((s, score, area))

    # 按标题得分排序，得分最高的当标题
    scored.sort(key=lambda x: (-x[1], -x[2]))
    title_shape = scored[0][0]

    # 剩下的按面积排序，最大的当内容
    remaining = [(s, area) for s, _, area in scored[1:]]
    if remaining:
        remaining.sort(key=lambda x: -x[1])
        content_shape = remaining[0][0]
    else:
        content_shape = title_shape

    return title_shape, content_shape


def _fill_image_slide(slide, title, extra):
    """全幅图片页：标题 + 居中图片 + 说明文字。"""
    from pptx.util import Inches
    image_path = extra.get("image", "")
    caption = extra.get("caption", "")

    # 找一个大的占位形状放图片
    shapes = list(text_shapes(slide))
    slots = [sh for sh in shapes if is_placeholder(sh.text_frame.text)]

    if slots and image_path:
        # 最大的形状放图片
        img_slot = max(slots, key=lambda s: (s.width or 0) * (s.height or 0))
        try:
            # 在形状位置嵌入图片
            left = img_slot.left
            top = img_slot.top
            width = img_slot.width
            height = img_slot.height
            # 清空占位文字
            set_text_keep_style(img_slot, "")
            # 嵌入图片
            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as e:
            set_text_keep_style(img_slot, f"【图片：{image_path}】{caption}")

    # 填标题
    title_shapes = [sh for sh in shapes if sh != img_slot if sh in slots]
    if title_shapes:
        set_text_keep_style(title_shapes[0], title)

    # 填说明文字
    if caption:
        remaining = [sh for sh in slots if sh not in (img_slot, title_shapes[0] if title_shapes else None)]
        if remaining:
            set_text_keep_style(remaining[0], caption)

    clear_decor(slide)
    clean_watermarks(slide)


def _fill_chart_slide(slide, title, extra):
    """全幅图表页：用 python-pptx 的 add_chart 嵌入图表。"""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_type_str = extra.get("chart_type", "bar")
    chart_data = extra.get("chart_data", {})
    categories = chart_data.get("categories", [])
    values = chart_data.get("values", [])

    if not categories or not values:
        # 数据不足，降级为文字
        fill_content(slide, title, [f"【图表数据缺失】{chart_type_str}: {categories}"], "text_only")
        return

    # 构建图表数据
    data = CategoryChartData()
    data.categories = categories
    data.add_series("数据", values)

    # 图表类型映射
    chart_types = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "bar_h": XL_CHART_TYPE.BAR_CLUSTERED,
    }
    ct = chart_types.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # 找一个大的形状位置放图表
    shapes = list(text_shapes(slide))
    slots = [sh for sh in shapes if is_placeholder(sh.text_frame.text)]

    if slots:
        chart_slot = max(slots, key=lambda s: (s.width or 0) * (s.height or 0))
        left = chart_slot.left
        top = chart_slot.top
        width = chart_slot.width
        height = chart_slot.height
        set_text_keep_style(chart_slot, "")
        slide.shapes.add_chart(ct, left, top, width, height, data)

    # 填标题
    title_shapes = [sh for sh in shapes if sh != chart_slot if sh in slots]
    if title_shapes:
        set_text_keep_style(title_shapes[0], title)

    clear_decor(slide)
    clean_watermarks(slide)


def _fill_table_slide(slide, title, extra):
    """全幅表格页：嵌入表格。"""
    from pptx.util import Inches

    table_data = extra.get("table_data", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers or not rows:
        fill_content(slide, title, [f"【表格数据缺失】"], "text_only")
        return

    # 找一个大的形状位置放表格
    shapes = list(text_shapes(slide))
    slots = [sh for sh in shapes if is_placeholder(sh.text_frame.text)]

    if slots:
        table_slot = max(slots, key=lambda s: (s.width or 0) * (s.height or 0))
        left = table_slot.left
        top = table_slot.top
        width = table_slot.width
        height = table_slot.height
        set_text_keep_style(table_slot, "")

        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)
        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        # 写表头
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
        # 写数据行
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)

    # 填标题
    title_shapes = [sh for sh in shapes if sh != table_slot if sh in slots]
    if title_shapes:
        set_text_keep_style(title_shapes[0], title)

    clear_decor(slide)
    clean_watermarks(slide)


def _fill_big_number_slide(slide, title, extra):
    """大数字页：页面中间放一个大数字，下面小字说明。"""
    number = extra.get("number", "0%")
    label = extra.get("label", "")

    shapes = list(text_shapes(slide))
    slots = [sh for sh in shapes if is_placeholder(sh.text_frame.text)]

    if not slots:
        clean_watermarks(slide)
        return

    # 找最大的形状放大数字
    big_slot = max(slots, key=lambda s: (s.width or 0) * (s.height or 0))

    # 写大数字（尽量大）
    set_text_keep_style(big_slot, number)
    # 设置字号尽可能大
    if big_slot.has_text_frame:
        from pptx.util import Pt
        for para in big_slot.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(72)  # 72pt 大数字
                run.font.bold = True
                from pptx.oxml.ns import qn
                rpr = run._r.get_or_add_rPr()
                ea = rpr.find(qn("a:ea"))
                if ea is None:
                    from lxml import etree
                    ea = etree.SubElement(rpr, qn("a:ea"))
                ea.set("typeface", "微软雅黑")

    # 找一个较小的形状放说明文字
    other_slots = [s for s in slots if s != big_slot]
    if other_slots and label:
        small_slot = min(other_slots, key=lambda s: (s.width or 0) * (s.height or 0))
        set_text_keep_style(small_slot, label)

    # 填标题
    title_shapes = [sh for sh in shapes if sh not in (big_slot, small_slot if other_slots else None) if sh in slots]
    if title_shapes:
        set_text_keep_style(title_shapes[0], title)

    clear_decor(slide)
    clean_watermarks(slide)


def _fill_comparison_slide(slide, title, extra):
    """对比页：左右或上下对比展示。"""
    items = extra.get("items", [])

    shapes = list(text_shapes(slide))
    slots = [sh for sh in shapes if is_placeholder(sh.text_frame.text)]

    if not slots or not items:
        clean_watermarks(slide)
        return

    # 将 slots 按面积排序，最大的放第一个对比项
    slots_by_area = sorted(slots, key=lambda s: (s.width or 0) * (s.height or 0), reverse=True)

    # 写对比内容
    for i, item in enumerate(items[:len(slots_by_area)]):
        set_text_keep_style(slots_by_area[i], item)

    # 填标题
    if len(slots_by_area) > len(items):
        set_text_keep_style(slots_by_area[len(items)], title)

    clear_decor(slide)
    clean_watermarks(slide)


def _fill_timeline_slide(slide, title, extra):
    """时间线页：按时间顺序展示。"""
    items = extra.get("items", [])

    shapes = list(text_shapes(slide))
    slots = [sh for sh in shapes if is_placeholder(sh.text_frame.text)]

    if not slots or not items:
        clean_watermarks(slide)
        return

    # 将 slots 按位置排序（从上到下、从左到右）
    slots_by_pos = sorted(slots, key=lambda s: (s.top or 0, s.left or 0))

    # 写时间线内容
    for i, item in enumerate(items[:len(slots_by_pos)]):
        set_text_keep_style(slots_by_pos[i], f"{i+1}. {item}")

    # 填标题
    if len(slots_by_pos) > len(items):
        set_text_keep_style(slots_by_pos[len(items)], title)

    clear_decor(slide)
    clean_watermarks(slide)


def drop_slides(prs, keep_indices):
    """删除未保留的页（1-based keep_indices），同时清理 sldIdLst 与关系。"""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    for i in range(len(ids) - 1, -1, -1):
        if (i + 1) not in keep_indices:
            rId = ids[i].get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
            prs.part.drop_rel(rId)
            sldIdLst.remove(ids[i])


def reorder_slides(prs, order):
    """按 order（slide 对象列表）重排 sldIdLst，使输出顺序为 封面→目录→各章→致谢。"""
    sldIdLst = prs.slides._sldIdLst
    id_by_slide = {}
    for sid in list(sldIdLst):
        rId = sid.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        id_by_slide[prs.part.related_part(rId).partname] = sid
    for slide in order:
        sid = id_by_slide.get(slide.part.partname)
        if sid is not None:
            sldIdLst.remove(sid)
            sldIdLst.append(sid)


def duplicate_slide(prs, src_slide, pristine_xml=None):
    """复制一页（含所有形状与样式），追加到末尾并返回新页。用于内容页复用。

    pristine_xml：源页"未被填写前"的 spTree 副本。必须传，否则复制到的是已填过
    上一章内容的页面（同一模板页被多章复用时会串内容）。
    """
    from copy import deepcopy
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    for sh in list(new_slide.shapes):      # 清掉版式带来的占位符
        sh._element.getparent().remove(sh._element)
    src_tree = pristine_xml if pristine_xml is not None else src_slide.shapes._spTree
    for child in src_tree:
        tag = child.tag.split('}')[-1]
        if tag in ("nvGrpSpPr", "grpSpPr"):   # 组属性新页已有，跳过
            continue
        new_slide.shapes._spTree.append(deepcopy(child))
    for rel in src_slide.part.rels.values():
        if rel.is_external:
            new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_slide.part.rels.get_or_add(rel.reltype, rel._target)
    return new_slide


def polish_slides(slides):
    """统一字体、段间距、清残余装饰，让填完文字的 PPT 视觉一致。

    只处理空字体的 run，不强制覆盖模板已有的字体选择。
    """
    from pptx.util import Pt
    from pptx.oxml.ns import qn

    for slide in slides:
        for sh in iter_text_shapes(slide):
            if not sh.has_text_frame:
                continue
            tf = sh.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    # 只处理没有字体设置的 run（保留模板原有字体）
                    if run.font.name is None:
                        run.font.name = "Calibri"
                        rpr = run._r.get_or_add_rPr()
                        rFonts = rpr.find(qn("a:latin"))
                        if rFonts is None:
                            from lxml import etree
                            rFonts = etree.SubElement(rpr, qn("a:latin"))
                        rFonts.set("typeface", "Calibri")
                        ea = rpr.find(qn("a:ea"))
                        if ea is None:
                            from lxml import etree
                            ea = etree.SubElement(rpr, qn("a:ea"))
                        ea.set("typeface", "宋体")

                # 段间距：段前 6pt、段后 2pt
                pPr = para._p.get_or_add_pPr()
                spc_before = pPr.find(qn("a:spcBef"))
                if spc_before is None:
                    from lxml import etree
                    spc_before = etree.SubElement(pPr, qn("a:spcBef"))
                spc_pct = spc_before.find(qn("a:spcPts"))
                if spc_pct is None:
                    from lxml import etree
                    for child in list(spc_before):
                        spc_before.remove(child)
                    spc_pct = etree.SubElement(spc_before, qn("a:spcPts"))
                spc_pct.set("val", str(int(6 * 100)))  # 6pt

                spc_after = pPr.find(qn("a:spcAft"))
                if spc_after is None:
                    from lxml import etree
                    spc_after = etree.SubElement(pPr, qn("a:spcAft"))
                spc_aft_pts = spc_after.find(qn("a:spcPts"))
                if spc_aft_pts is None:
                    from lxml import etree
                    for child in list(spc_after):
                        spc_after.remove(child)
                    spc_aft_pts = etree.SubElement(spc_after, qn("a:spcPts"))
                spc_aft_pts.set("val", str(int(2 * 100)))  # 2pt

    print("  ✓ 已统一段间距(段前6pt/段后2pt)，保留模板原有字体")


def build(template: Path, content_path: Path, output: Path):
    data = json.loads(content_path.read_text(encoding="utf-8"))
    key = next((k for k in TEMPLATE_MAP if k in template.stem), None)
    if key is None:
        raise SystemExit(f"❌ 未识别的模板：{template.name}（应含 模板1/模板2/模板3）")
    roles = TEMPLATE_MAP[key]

    # 自动修复内容密度问题（合并空洞页、精简超长 bullets）
    data = auto_fix_content_density(data)

    prs = Presentation(str(template))
    slides = list(prs.slides)
    chapters = data.get("chapters", [])
    if not chapters:
        raise SystemExit("❌ content.json 缺少 chapters（章节列表）")

    cover_slide = slides[roles["cover"] - 1]
    toc_slide = slides[roles["toc"] - 1]
    thanks_slide = slides[roles["thanks"] - 1] if roles.get("thanks") else None

    fill_cover(cover_slide, data)
    fill_toc(toc_slide, [c["name"] for c in chapters])

    order = [cover_slide, toc_slide]
    sec_pool = [slides[i - 1] for i in roles["section"]]
    con_pool = [slides[i - 1] for i in roles["content"]]

    # 关键：在任何填写之前先留存"干净"的形状树副本，供复制页使用，
    # 否则复用同一模板页时会复制到上一章已填的内容。
    from copy import deepcopy
    pristine = {s.part.partname: deepcopy(s.shapes._spTree)
                for s in sec_pool + con_pool}

    for ci, ch in enumerate(chapters):
        # 章节过渡页：模板页够用就直接用，不够则复制最后一个章节页样式
        if ci < len(sec_pool):
            sec = sec_pool[ci]
        else:
            src = sec_pool[-1]
            sec = duplicate_slide(prs, src, pristine.get(src.part.partname))
        fill_section(sec, ci + 1, ch["name"])
        order.append(sec)

        # 内容页：轮转取模板内容页样式；首次直接用、其后按干净副本复制
        for si, spec in enumerate(ch.get("slides", [])):
            base = con_pool[si % len(con_pool)]
            if base in order:            # 该模板页已用过 → 复制一份干净的
                slide = duplicate_slide(prs, base, pristine.get(base.part.partname))
            else:
                slide = base
            layout = spec.get("layout", "text_only")
            extra = spec.get("extra", {})
            if layout != "text_only" and extra:
                fill_content(slide, spec.get("title", ""), spec.get("bullets", []), layout, extra)
            else:
                fill_content(slide, spec.get("title", ""), spec.get("bullets", []))
            # 写入演讲者备注
            notes = spec.get("notes", "")
            if notes:
                set_slide_notes(slide, notes)
            order.append(slide)

    if thanks_slide is not None:
        clear_decor(thanks_slide)
        clean_watermarks(thanks_slide)
        cover = data.get("cover", {})
        for sh in iter_text_shapes(thanks_slide):
            t = sh.text_frame.text.strip()
            if t in ("汇报人", "答辩人", "汇报人：", "答辩人：") and cover.get("作者姓名"):
                set_text_keep_style(sh, f"{t.rstrip('：')}：{cover['作者姓名']}")
            elif t in ("指导老师", "指导老师：", "指导教师", "指导教师：") and cover.get("指导教师"):
                set_text_keep_style(sh, f"{t.rstrip('：')}：{cover['指导教师']}")
        order.append(thanks_slide)

    # 删除所有未用到的模板页，再按逻辑顺序重排
    keep_partnames = {s.part.partname for s in order}
    keep_idx = [i + 1 for i, s in enumerate(prs.slides) if s.part.partname in keep_partnames]
    drop_slides(prs, keep_idx)
    reorder_slides(prs, order)

    # ── 美化：统一字体 + 段间距 + 清残余装饰 ──
    polish_slides(order)

    # 清掉内置模板带来的第三方个人信息：模板是借用的往届答辩稿，
    # docProps/core.xml 里留着原作者姓名与邮箱。不清则每份汇报 PPT 都会在
    # 属性面板显示陌生人的邮箱，既泄露他人隐私，也像是别人做的稿子。
    cp = prs.core_properties
    scrubbed = [k for k, v in (("author", cp.author), ("last_modified_by", cp.last_modified_by))
                if v]
    au = str(data.get("author", "") or "")
    if any(m in au for m in ("〈", "〉", "XXX", "待补", "待定")):
        au = ""
    cp.author = au
    cp.last_modified_by = au
    cp.title = str(data.get("title", "") or "")

    prs.save(str(output))
    print("saved:", output)
    validate(output, data, scrubbed=scrubbed)
    validate_content_density(data)
    # 调用 pptx 技能的 validate.py 做文件级验证（XML schema / relationships / charts）
    _run_pptx_validate(output, template)
    # 生成视觉预览图片
    _generate_preview_images(output)


def validate(path, data, scrubbed=None):
    prs = Presentation(str(path))
    all_text = []
    for s in prs.slides:
        for sh in iter_text_shapes(s):
            all_text.append(sh.text_frame.text)
    text = "\n".join(all_text)
    leftover_ph = [h for h in ("请输入", "点击此处", "在这里输入", "JUNE 12th", "ADD YOUR") if h in text]
    leftover_wm = [w for w in WATERMARKS if w in text]
    print("=" * 50)
    print("自检：")
    print(f"  页数: {len(prs.slides)}")
    print(f"  题目已写入: {data.get('title','')[:20] in text}")
    print(f"  章节数: {len(data.get('chapters', []))}")
    print(f"  水印已清: {not leftover_wm} {leftover_wm if leftover_wm else ''}")
    print(f"  残留占位: {leftover_ph if leftover_ph else '无'}")
    if leftover_ph:
        print("  ↑ 提示：复杂图表页可能仍有占位文字/图片，需在 PowerPoint 手动微调")
    if scrubbed:
        print(f"  文档属性已清理: {scrubbed}（模板原作者姓名/邮箱已清空）")


def _run_pptx_validate(output_path, template_path):
    """调用 pptx 技能的 validate.py 做文件级验证。"""
    import subprocess
    # pptx 技能的 validate.py 路径
    pptx_skill_dir = Path.home() / ".claude" / "skills" / "pptx"
    validate_script = pptx_skill_dir / "scripts" / "office" / "validate.py"
    if not validate_script.exists():
        return  # pptx 技能未安装，跳过
    try:
        result = subprocess.run(
            ["python3", str(validate_script), str(output_path), "--original", str(template_path)],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode != 0:
            print(f"  ⚠ pptx validate 发现问题：")
            for line in result.stdout.strip().split("\n")[:10]:  # 最多显示10行
                print(f"    {line}")
        else:
            print("  ✓ pptx validate 通过")
    except (subprocess.TimeoutExpired, FileNotFoundError):
        pass  # 静默跳过，不影响主流程


def _generate_preview_images(pptx_path):
    """将 PPT 转为图片预览（需要 LibreOffice）。"""
    import subprocess
    output_dir = pptx_path.parent
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"

    # 检查 pptx 技能的 soffice.py 是否存在
    pptx_skill_dir = Path.home() / ".claude" / "skills" / "pptx"
    soffice_script = pptx_skill_dir / "scripts" / "office" / "soffice.py"
    if not soffice_script.exists():
        return  # pptx 技能未安装，跳过

    try:
        # 先转 PDF
        result = subprocess.run(
            ["python3", str(soffice_script), "--headless", "--convert-to", "pdf", str(pptx_path)],
            capture_output=True, text=True, timeout=60,
            cwd=str(output_dir)
        )
        if result.returncode != 0 or not pdf_path.exists():
            return

        # 再用 pdftoppm 转图片
        import glob
        old_images = glob.glob(str(output_dir / "slide-*.jpg"))
        for img in old_images:
            os.remove(img)

        result = subprocess.run(
            ["pdftoppm", "-jpeg", "-r", "150", str(pdf_path), str(output_dir / "slide")],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode == 0:
            preview_files = sorted(glob.glob(str(output_dir / "slide-*.jpg")))
            if preview_files:
                print(f"  📸 视觉预览已生成：{len(preview_files)} 张图片")
                print(f"  📂 路径：{output_dir}/slide-*.jpg")
                # 列出前几张
                for f in preview_files[:3]:
                    print(f"    {Path(f).name}")
                if len(preview_files) > 3:
                    print(f"    ... 共 {len(preview_files)} 张")
    except (subprocess.TimeoutExpired, FileNotFoundError):
        pass  # 静默跳过，不影响主流程


def validate_content_density(data):
    """检查 ppt_content.json 的内容密度是否达标，给出改进建议。"""
    chapters = data.get("chapters", [])
    total_slides = sum(len(ch.get("slides", [])) for ch in chapters)
    issues = []

    # 总页数检查（含封面/目录/致谢约 +3 页）
    effective = total_slides + 3
    if effective > 20:
        issues.append(f"总页数 {effective} 页偏多，8 分钟汇报建议 12–15 页；"
                      "宁可每页说透，不要拆太多页导致每页太空")
    if effective < 10:
        issues.append(f"总页数仅 {effective} 页，建议至少 12 页以保证内容完整")

    # 逐章检查：每页 bullets 数量
    thin_pages = []
    for ch in chapters:
        for sl in ch.get("slides", []):
            n_bullets = len(sl.get("bullets", []))
            if n_bullets < 3:
                thin_pages.append(f"「{ch['name']}」→「{sl.get('title', '')[:20]}」仅 {n_bullets} 条")
    if thin_pages:
        issues.append("以下页面内容过薄（<3 条 bullets），建议补充或合并到相邻页：" + "；".join(thin_pages[:5]))

    # 检查是否有图表占位
    has_diagram = False
    for ch in chapters:
        for sl in ch.get("slides", []):
            for b in sl.get("bullets", []):
                if "【图" in b or "【表" in b or "路线图" in b:
                    has_diagram = True
    if not has_diagram:
        issues.append("未发现图表占位（【图：xxx】/【表：xxx】），研究框架/路线图/时间表建议嵌图")

    # 检查章节标题是否太泛
    generic_titles = {"研究背景", "文献综述", "研究方法", "创新之处", "实施计划",
                      "选题意义", "研究框架", "研究思路"}
    for ch in chapters:
        for sl in ch.get("slides", []):
            if sl.get("title", "") in generic_titles:
                issues.append(f"页面标题「{sl['title']}」太泛——应写具体观点，不是章节名")
                break

    # 检查占位残留
    for ch in chapters:
        for sl in ch.get("slides", []):
            for b in sl.get("bullets", []):
                if "见下页" in b or "添加" in b:
                    issues.append(f"「{ch['name']}」存在占位残留：「{b}」— 要么嵌入内容，要么写【图：xxx】占位")

    if issues:
        print("\n⚠ 内容密度自检：")
        for i, issue in enumerate(issues, 1):
            print(f"  {i}. {issue}")
        print("  ↑ 以上问题不影响生成，但会影响汇报质量。建议补充后重新生成。")
    else:
        print("\n✓ 内容密度自检通过")


def auto_fix_content_density(data):
    """自动修复内容密度问题：合并空洞页、精简超长 bullets。"""
    chapters = data.get("chapters", [])
    fixes = []

    for ch in chapters:
        slides = ch.get("slides", [])
        if not slides:
            continue

        # 修复 1：合并空洞页（bullets < 2 的页合并到前一页）
        merged_slides = []
        for sl in slides:
            bullets = sl.get("bullets", [])
            if len(bullets) < 2 and merged_slides:
                # 合并到前一页
                prev = merged_slides[-1]
                prev_bullets = prev.get("bullets", [])
                prev_bullets.extend(bullets)
                prev["bullets"] = prev_bullets
                fixes.append(f"合并空洞页「{sl.get('title', '')}」到「{prev.get('title', '')}」")
            else:
                merged_slides.append(sl)
        ch["slides"] = merged_slides

        # 修复 2：精简超长 bullets（超过 25 字的截断）
        for sl in ch.get("slides", []):
            bullets = sl.get("bullets", [])
            new_bullets = []
            for b in bullets:
                if len(b) > 25:
                    # 按标点截断
                    parts = re.split(r"[，。；、]", b)
                    shortened = parts[0] if len(parts[0]) > 10 else b[:25]
                    new_bullets.append(shortened)
                    fixes.append(f"精简超长 bullet：「{b[:30]}...」→「{shortened}」")
                else:
                    new_bullets.append(b)
            sl["bullets"] = new_bullets

    if fixes:
        print(f"\n🔧 自动修复了 {len(fixes)} 个问题：")
        for fix in fixes[:10]:  # 最多显示 10 个
            print(f"  • {fix}")
    else:
        print("\n✓ 无需自动修复")

    return data


def main():
    ap = argparse.ArgumentParser(description="用内置北航模板生成开题汇报 PPT")
    ap.add_argument("--template", required=True, type=Path)
    ap.add_argument("--content", required=True, type=Path)
    ap.add_argument("--output", required=True, type=Path)
    a = ap.parse_args()
    build(a.template, a.content, a.output)


if __name__ == "__main__":
    main()
